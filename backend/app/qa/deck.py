"""Renders a QA report as a slide deck and as Markdown.

Markdown is the guaranteed output; the deck is a bonus. python-pptx is an
optional dependency (requirements-qa.txt), so if it is missing the report still
gets written and *says* the deck was skipped and how to enable it. A status
report that silently loses half of itself is worse than one that admits it.

The two renderers deliberately carry the same content: whoever reads the deck in
a review and whoever greps the Markdown in a terminal must not end up with
different beliefs about what passed.
"""
from __future__ import annotations

from pathlib import Path
from types import SimpleNamespace
from typing import Optional

DECK_HINT = ("python-pptx is not installed, so only the Markdown report was "
             "written. Enable the deck with `pip install -r "
             "backend/requirements-qa.txt`, or add `--with python-pptx` to the "
             "`uv run` invocation.")

STATUS_MARK = {"pass": "PASS", "fail": "FAIL", "skip": "SKIP"}
STATUS_RGB = {"pass": (0x1B, 0x5E, 0x20), "fail": (0xB7, 0x1C, 0x1C),
              "skip": (0x61, 0x61, 0x61)}


def _pptx_api() -> Optional[SimpleNamespace]:
    """The handful of python-pptx names this module uses, or None.

    Bundled into one lookup so the absent-dependency path has a single point of
    control — which is also the seam the tests use to exercise it.
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
    except ImportError:
        return None
    return SimpleNamespace(Presentation=Presentation, RGBColor=RGBColor,
                           Inches=Inches, Pt=Pt)


# --------------------------------------------------------------------- markdown

def _headline(report: dict) -> str:
    if report["status"] == "running":
        done = len(report.get("flows", []))
        return (f"Running — {done}/{report.get('flows_total', '?')} flows, "
                f"{report['passed']}/{report['total']} checks so far")
    verdict = "all checks passed" if report["status"] == "done" else "failures present"
    return f"**{report['passed']} / {report['total']} checks passed** — {verdict}"


def render_markdown(report: dict) -> str:
    """The whole report as one document. Screenshot links are relative, so the
    file reads correctly both on disk and when served from its own directory."""
    flows = report.get("flows", [])
    lines = [
        "# Application status — CV Dataset Explorer",
        "",
        _headline(report),
        "",
        f"- Run `{report['run_id']}` · started {report.get('started_at')}"
        + (f" · finished {report['finished_at']}" if report.get("finished_at") else "")
        + (f" · {report['duration_s']:.0f}s" if report.get("duration_s") else ""),
        f"- Target {report.get('base_url')} (API {report.get('api_url')})",
        f"- {report.get('flows_passed', 0)}/{report.get('flows_total', len(flows))} "
        "workflows fully green",
    ]
    if report.get("error"):
        lines.append(f"- **Run error:** {report['error']}")
    if report.get("deck_note"):
        lines.append(f"- {report['deck_note']}")
    lines += ["", "| Workflow | Status | Checks | Time |", "|---|---|---|---|"]
    for f in flows:
        passed = sum(1 for c in f["checks"] if c["ok"])
        lines.append(f"| {f['name']} | {STATUS_MARK.get(f['status'], f['status'])} "
                     f"| {passed}/{len(f['checks'])} | {f.get('duration_s') or 0:.1f}s |")
    for name in report.get("pending") or []:
        lines.append(f"| {name} | pending | – | – |")

    for f in flows:
        passed = sum(1 for c in f["checks"] if c["ok"])
        lines += ["", f"## {f['name']} — {STATUS_MARK.get(f['status'], f['status'])} "
                      f"({passed}/{len(f['checks'])}, {f.get('duration_s') or 0:.1f}s)", ""]
        if f.get("detail"):
            lines += [f"> {f['detail']}", ""]
        if f.get("screenshot"):
            lines += [f"![{f['name']}]({Path(f['screenshot']).name})", ""]
        for c in f["checks"]:
            detail = f" — {c['detail']}" if c.get("detail") else ""
            lines.append(f"- {'PASS' if c['ok'] else 'FAIL'} {c['name']}{detail}")

    for heading, key in (("Console errors", "console_errors"),
                         ("Failed network requests", "network_failures")):
        lines += ["", f"## {heading}", ""]
        entries = report.get(key) or []
        lines += [f"- `{e}`" for e in entries] or ["- none"]
    return "\n".join(lines) + "\n"


# ------------------------------------------------------------------------ deck

def render_deck(report: dict, path: Path, media_dir: Path) -> Optional[Path]:
    """Write a .pptx: title slide, then one slide per flow. None if python-pptx
    is unavailable."""
    api = _pptx_api()
    if api is None:
        return None
    Inches, Pt, RGB = api.Inches, api.Pt, api.RGBColor

    prs = api.Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    def textbox(slide, x, y, w, h):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        box.text_frame.word_wrap = True
        return box.text_frame

    def write(frame, text, *, size=14, bold=False, rgb=(0x21, 0x21, 0x21), first=False):
        para = frame.paragraphs[0] if first else frame.add_paragraph()
        run = para.add_run()
        run.text = text
        run.font.size, run.font.bold = Pt(size), bold
        run.font.color.rgb = RGB(*rgb)
        return para

    # ---- title slide
    s = prs.slides.add_slide(blank)
    tf = textbox(s, 0.9, 1.7, 11.5, 3.4)
    write(tf, "Application status", size=40, bold=True, first=True)
    write(tf, "CV Dataset Explorer — autonomous UI sweep", size=20,
          rgb=(0x55, 0x55, 0x55))
    verdict = ("ALL CHECKS PASSED" if report["status"] == "done"
               else "FAILURES PRESENT" if report["status"] == "failed"
               else "RUN IN PROGRESS")
    write(tf, f"{report['passed']} / {report['total']} checks · "
              f"{report.get('flows_passed', 0)} / "
              f"{report.get('flows_total', 0)} workflows green · {verdict}",
          size=22, bold=True,
          rgb=STATUS_RGB["pass" if report["status"] == "done" else "fail"])
    meta = [f"Run {report['run_id']} · {report.get('finished_at') or report.get('started_at')}",
            f"Target {report.get('base_url')} · driven through Chrome by Playwright"]
    if report.get("console_errors"):
        meta.append(f"{len(report['console_errors'])} console error(s) — see the last slide")
    if report.get("error"):
        meta.append(f"Run error: {report['error']}")
    for line in meta:
        write(tf, line, size=13, rgb=(0x55, 0x55, 0x55))

    # ---- one slide per flow
    for f in report.get("flows", []):
        s = prs.slides.add_slide(blank)
        passed = sum(1 for c in f["checks"] if c["ok"])
        head = textbox(s, 0.5, 0.35, 12.3, 0.9)
        write(head, f['name'], size=26, bold=True, first=True)
        write(head, f"{STATUS_MARK.get(f['status'], f['status'])} · {passed}/"
                    f"{len(f['checks'])} checks · {f.get('duration_s') or 0:.1f}s"
                    + (f" · {f['detail']}" if f.get("detail") else ""),
              size=13, bold=True, rgb=STATUS_RGB.get(f["status"], (0x21, 0x21, 0x21)))

        body = textbox(s, 0.5, 1.6, 5.2, 5.4)
        for i, c in enumerate(f["checks"]):
            detail = f"  ({c['detail']})" if c.get("detail") else ""
            write(body, f"{'PASS' if c['ok'] else 'FAIL'}  {c['name']}{detail}"[:180],
                  size=11, bold=not c["ok"],
                  rgb=STATUS_RGB["pass" if c["ok"] else "fail"], first=(i == 0))
        if not f["checks"]:
            write(body, "no checks recorded", size=11, rgb=STATUS_RGB["skip"], first=True)

        shot = media_dir / Path(f["screenshot"]).name if f.get("screenshot") else None
        if shot and shot.exists():
            # Width-fitted at the sweep's 16:10 viewport aspect; height follows.
            s.shapes.add_picture(str(shot), Inches(6.0), Inches(1.7), width=Inches(6.8))

    # ---- console errors, if any
    problems = (report.get("console_errors") or []) + (report.get("network_failures") or [])
    if problems:
        s = prs.slides.add_slide(blank)
        tf = textbox(s, 0.5, 0.5, 12.3, 6.5)
        write(tf, "Console errors and failed requests", size=26, bold=True, first=True)
        for line in problems[:30]:
            write(tf, line[:200], size=12, rgb=STATUS_RGB["fail"])
        if len(problems) > 30:
            write(tf, f"… and {len(problems) - 30} more", size=12, rgb=STATUS_RGB["skip"])

    prs.save(str(path))
    return path


# ----------------------------------------------------------------------- both

def write(report: dict, out_dir: Path, *, make_deck: bool = True) -> dict:
    """Render both documents into `out_dir`.

    Returns file *names* rather than URLs: how artifacts are addressed is the
    runner's business, and keeping it there is what stops this module from
    importing its own caller.
    """
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    result = {"markdown_file": "report.md", "deck_file": None, "deck_note": None}

    # Deck first, so a skipped or failed deck is *in* the Markdown rather than
    # only in the caller's return value — the surviving document has to say what
    # is missing from it.
    if not make_deck:
        result["deck_note"] = "Deck skipped (--no-deck)."
    else:
        try:
            path = render_deck(report, out_dir / "report.pptx", out_dir)
            if path is None:
                result["deck_note"] = DECK_HINT
            else:
                result["deck_file"] = path.name
        except Exception as e:                        # noqa: BLE001
            # A deck that fails to render must not lose the run that produced it.
            result["deck_note"] = (f"Deck render failed ({type(e).__name__}: {e}); "
                                   "the Markdown report is complete.")

    (out_dir / "report.md").write_text(
        render_markdown({**report, "deck_note": result["deck_note"]}))
    return result
